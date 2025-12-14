"""
导出路由
"""

import sys
import base64
import tempfile
from pathlib import Path
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse

# 添加项目根目录到 Python 路径
project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from src.exporter import PDFExporter
from ..models import ExportRequest

router = APIRouter(prefix="/api", tags=["export"])


@router.post("/export")
async def export_presentation(request: ExportRequest):
    """
    导出演示文稿
    
    Args:
        request: 导出请求，包含所有幻灯片和格式
        
    Returns:
        FileResponse: 导出的文件
    """
    try:
        # 创建临时目录保存图片
        temp_dir = Path(tempfile.mkdtemp())
        image_paths = []
        
        # 解码并保存所有图片
        for idx, slide in enumerate(request.slides):
            image_data = base64.b64decode(slide.image_base64)
            image_path = temp_dir / f"slide_{idx + 1}.png"
            
            with open(image_path, 'wb') as f:
                f.write(image_data)
            
            image_paths.append(str(image_path))
        
        # 根据格式导出
        if request.format == "pdf":
            output_path = temp_dir / "presentation.pdf"
            exporter = PDFExporter()
            exporter.export(image_paths, str(output_path))
            
            return FileResponse(
                path=str(output_path),
                media_type="application/pdf",
                filename="presentation.pdf"
            )
        
        elif request.format == "pptx":
            output_path = temp_dir / "presentation.pptx"
            _export_pptx(image_paths, str(output_path))
            
            return FileResponse(
                path=str(output_path),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                filename="presentation.pptx"
            )
        
        else:
            raise HTTPException(
                status_code=400,
                detail=f"不支持的导出格式: {request.format}"
            )
    
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"导出失败: {str(e)}"
        )


def _export_pptx(image_paths: list, output_path: str):
    """
    导出为 PPTX 格式
    
    Args:
        image_paths: 图片路径列表
        output_path: 输出路径
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise Exception("需要安装 python-pptx: pip install python-pptx")
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸为 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 添加每一页
    for image_path in image_paths:
        # 使用空白布局
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 添加图片，填充整个幻灯片
        slide.shapes.add_picture(
            image_path,
            0, 0,
            width=prs.slide_width,
            height=prs.slide_height
        )
    
    # 保存
    prs.save(output_path)
